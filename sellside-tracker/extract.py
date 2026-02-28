#!/usr/bin/env python3
"""
Sellside Report Tracker — Phase 0 极简版
PDF → AI 结构化提取 → YAML 存储 → QoQ 对比

Usage:
    python extract.py "path/to/report.pdf" [--ticker GOOG] [--quarter Q4-2025]
    python extract.py --history GOOG
"""

import argparse
import csv
import json
import os
import re
import sqlite3
import sys
from datetime import datetime
from pathlib import Path

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

import pymupdf
import yaml

# ── Paths ────────────────────────────────────────────────────────────────────

SKILL_DIR = Path(__file__).parent
TEMPLATES_DIR = SKILL_DIR / "templates"
DATA_DIR = SKILL_DIR / "data"
DATA_DIR.mkdir(exist_ok=True)

sys.path.insert(0, str(Path.home() / ".claude" / "skills"))
from shared.obsidian_utils import create_note
from shared.dashboard_updater import update_dashboard

VAULT = Path.home() / "Documents" / "Obsidian Vault"


# ── File Extraction ──────────────────────────────────────────────────────────

def extract_pdf_text(pdf_path: str) -> str:
    """Extract full text from PDF with page markers."""
    doc = pymupdf.open(pdf_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        pages.append(f"[PAGE {i + 1}]\n{text}")
    return "\n\n".join(pages)


def extract_pptx_text(pptx_path: str) -> str:
    """Extract full text from PPTX with slide markers."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))
        slides.append(f"[PAGE {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_file_text(file_path: str) -> str:
    """Extract text from PDF or PPTX."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    else:
        print(f"Error: Unsupported file type: {ext}", file=sys.stderr)
        sys.exit(1)


# ── Template Loading ─────────────────────────────────────────────────────────

TICKER_ALIASES = {
    "GOOG": ["google", "alphabet"], "GOOGL": ["google", "alphabet"],
    "META": ["meta", "facebook"], "MSFT": ["microsoft"],
    "AAPL": ["apple"], "AMZN": ["amazon"],
    "NVDA": ["nvidia"], "TSLA": ["tesla"],
}


def load_template(ticker: str) -> list[dict]:
    """Load CSV template for a ticker. Tries ticker name then aliases."""
    candidates = [ticker.lower()]
    candidates.extend(TICKER_ALIASES.get(ticker.upper(), []))
    csv_path = None
    for name in candidates:
        p = TEMPLATES_DIR / f"{name}.csv"
        if p.exists():
            csv_path = p
            break
    if not csv_path:
        return None
    with open(csv_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        return [row for row in reader]


def template_to_prompt_fields(fields: list[dict]) -> str:
    """Convert template fields to prompt instruction text."""
    lines = []
    for f in fields:
        req = " [必填]" if f["required"] == "yes" else ""
        note = f" ({f['notes']})" if f.get("notes", "").strip() else ""
        lines.append(f"- {f['key']} ({f['type']}): {f['label']}{note}{req}")
    return "\n".join(lines)


# ── AI Extraction ────────────────────────────────────────────────────────────

def build_extraction_prompt(pdf_text: str, fields_text: str, ticker: str, quarter: str) -> str:
    return f"""你是金融数据提取助手。从以下卖方分析师季度报告中精确提取指标。

公司: {ticker}
季度: {quarter}

## 提取规则
1. 每个指标标注来源页码 source_page（整数）
2. 报告中未提及的指标设为 null，绝不猜测
3. 需要计算得出的值，标注 confidence: calculated
4. percent 类型保留原始格式如 "17%"
5. number 类型统一为百万美元(mn USD)，若原文是亿人民币请在 note 字段注明
6. text 类型用原文关键句，不超过 80 字
7. [必填] 标记的字段务必尽力提取

## 要提取的指标
{fields_text}

## 输出格式
严格输出 YAML，不要 markdown 代码块标记，格式如下：

ticker: {ticker}
quarter: {quarter}
report_date: YYYY-MM
source: "分析师/机构名"
metrics:
  <key>:
    value: <值>
    source_page: <页码>
    confidence: high|medium|calculated
    note: <可选备注>

## 报告全文
{pdf_text}
"""


def build_summary_prompt(pdf_text: str, ticker: str, quarter: str) -> str:
    return f"""你是一位资深买方分析师助理。请将以下卖方季度跟踪报告转化为结构化的投资笔记。

公司: {ticker}
季度: {quarter}

## 输出要求
用中文输出 Markdown 格式（不要代码块包裹），按以下结构组织：

### 核心观点
- 用 3-5 个 bullet 概括卖方本季度的核心判断（多空态度、关键变化、估值结论）

### Search 广告
- 收入增速、量价拆分（clicks vs CPC）
- AI 对搜索的影响（AI Overviews、AI Mode 渗透率、query 变化）
- 三方数据验证（如有 SimilarWeb、Yipit 等）
- 竞争格局（份额变化、ChatGPT/Perplexity 冲击）

### Cloud / GCP
- 收入增速、利润率变化
- 客户获取（新客户、大单、backlog）
- AI 云产品渗透率
- Capex 规模及投入方向

### YouTube & 订阅
- 广告收入增速
- Shorts 货币化、CTV 份额
- 订阅业务规模（Google One、Premium、TV）

### AI 生态
- Gemini 模型进展（MAU、tokens 处理量、新模型发布）
- AI agent 平台（AntiGravity、Genie 等）
- 内部 AI 效率提升（代码生成、运营自动化）

### 公司财务 & 资本配置
- 收入增速、OP 增速、利润率变化
- 费用结构（R&D、S&M 杠杆效应）
- Capex 指引、FCF 影响
- 回购/分红力度变化

### 估值
- 卖方的估值框架（整体 PE、SOTP 拆分）
- 各业务隐含估值（Search、Cloud、YouTube、Waymo）
- 卖方 vs 市场一致预期的差异（如有）

### Other Bets
- Waymo（融资、rides、城市扩展）
- DeepMind 进展

## 规则
1. 保留具体数字和百分比，标注来源页码如 (p.5)
2. 区分公司披露数据和卖方自己的推算/判断
3. 如果某个 section 报告中未涉及，写"本报告未涉及"
4. 不要添加报告中没有的信息

## 报告全文
{pdf_text}
"""


def build_generic_extraction_prompt(pdf_text: str, ticker: str) -> str:
    """Generic extraction prompt — no CSV template needed."""
    return f"""You are a financial data extraction assistant. Extract the following standardized fields from this sell-side research report.

Company: {ticker}

## Extract these fields (JSON format):

{{
    "ticker": "{ticker}",
    "date": "YYYY-MM-DD (report date)",
    "firm": "Brokerage firm name",
    "analyst": "Lead analyst name",
    "rating": "Buy/Overweight/Hold/Neutral/Sell/Underweight (exact rating)",
    "prior_rating": "Previous rating if mentioned, else null",
    "target_price": numeric value (no currency symbol),
    "prior_target": numeric value if mentioned, else null,
    "currency": "USD/HKD/CNY/EUR/GBP",
    "key_thesis": "1-3 sentence core thesis in English",
    "catalysts": ["list", "of", "upcoming", "catalysts"],
    "risks": ["list", "of", "key", "risks"],
    "estimates": {{
        "revenue_current_q": "$ value or null",
        "revenue_next_q": "$ value or null",
        "eps_current_q": "$ value or null",
        "eps_next_q": "$ value or null"
    }}
}}

## Rules:
1. Output ONLY valid JSON, no markdown fences
2. If a field is not found in the report, set it to null
3. For rating, normalize to: Buy, Overweight, Hold, Neutral, Underweight, Sell
4. target_price and prior_target should be numeric only (no $ or currency)
5. key_thesis should capture the analyst's main argument

## Report text:
{pdf_text}"""


def call_gemini(prompt: str) -> str:
    """Call Gemini API for extraction."""
    from google import genai

    # Load API key
    config_path = Path.home() / ".claude" / "skills" / "prompt-optimizer" / "data" / "config.json"
    if config_path.exists():
        cfg = json.loads(config_path.read_text(encoding="utf-8"))
        api_key = cfg.get("GEMINI_API_KEY", "")
    else:
        api_key = os.environ.get("GEMINI_API_KEY", "")

    if not api_key:
        print("Error: No Gemini API key found.", file=sys.stderr)
        sys.exit(1)

    client = genai.Client(api_key=api_key)
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=prompt,
        config=genai.types.GenerateContentConfig(
            temperature=0.1,
            max_output_tokens=4096,
        ),
    )
    return response.text


# ── Validation ───────────────────────────────────────────────────────────────

def validate_extraction(data: dict, fields: list[dict], prev_data: dict | None) -> list[str]:
    """Validate extracted data. Returns list of warnings."""
    warnings = []
    metrics = data.get("metrics", {})

    for f in fields:
        key = f["key"]
        m = metrics.get(key)

        # Required field missing
        if f["required"] == "yes" and (m is None or m.get("value") is None):
            warnings.append(f"⚠️ 必填字段缺失: {f['label']} ({key})")
            continue

        if m is None or m.get("value") is None:
            continue

        val = m["value"]

        # Range check for percents
        if f["type"] == "percent" and isinstance(val, str) and "%" in val:
            try:
                num = float(val.replace("%", "").strip())
                if abs(num) > 200:
                    warnings.append(f"⚠️ 异常值: {f['label']} = {val}")
            except ValueError:
                pass

        # QoQ deviation check
        if prev_data and key in prev_data.get("metrics", {}):
            prev_val = prev_data["metrics"][key].get("value")
            if prev_val and f["type"] == "percent":
                try:
                    cur = float(str(val).replace("%", "").strip())
                    prev = float(str(prev_val).replace("%", "").strip())
                    delta = abs(cur - prev)
                    if delta > 20:
                        warnings.append(
                            f"⚠️ 大幅变化: {f['label']} {prev_val} → {val} (Δ{delta:.0f}pp)"
                        )
                except ValueError:
                    pass

    return warnings


# ── Storage ──────────────────────────────────────────────────────────────────

def load_history(ticker: str) -> dict:
    """Load all historical quarterly data for a ticker."""
    path = DATA_DIR / f"{ticker.lower()}_quarterly.yaml"
    if path.exists():
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    return {}


def save_history(ticker: str, history: dict):
    """Save quarterly data for a ticker."""
    path = DATA_DIR / f"{ticker.lower()}_quarterly.yaml"
    with open(path, "w", encoding="utf-8") as f:
        yaml.dump(history, f, allow_unicode=True, default_flow_style=False, sort_keys=False)
    print(f"✅ 数据已保存: {path}")


def get_previous_quarter(history: dict, current_quarter: str) -> dict | None:
    """Get the most recent quarter before current."""
    quarters = sorted(history.keys())
    if current_quarter in quarters:
        idx = quarters.index(current_quarter)
        if idx > 0:
            return history[quarters[idx - 1]]
    elif quarters:
        return history[quarters[-1]]
    return None


# ── SQLite Views Storage ──────────────────────────────────────────────────────

def _get_views_db() -> sqlite3.Connection:
    db_path = DATA_DIR / "sellside_views.db"
    conn = sqlite3.connect(str(db_path))
    conn.execute("""CREATE TABLE IF NOT EXISTS views (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        ticker TEXT NOT NULL,
        date TEXT NOT NULL,
        firm TEXT NOT NULL,
        analyst TEXT,
        rating TEXT,
        prior_rating TEXT,
        target_price REAL,
        prior_target REAL,
        currency TEXT DEFAULT 'USD',
        key_thesis TEXT,
        catalysts TEXT,
        risks TEXT,
        estimates TEXT,
        source_file TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP,
        UNIQUE(ticker, date, firm)
    )""")
    conn.execute("CREATE INDEX IF NOT EXISTS idx_views_ticker ON views(ticker)")
    conn.execute("CREATE INDEX IF NOT EXISTS idx_views_firm ON views(ticker, firm)")
    conn.commit()
    return conn


def save_view(ticker: str, view: dict, source_file: str = ""):
    conn = _get_views_db()
    conn.execute("""INSERT OR REPLACE INTO views
        (ticker, date, firm, analyst, rating, prior_rating, target_price, prior_target,
         currency, key_thesis, catalysts, risks, estimates, source_file)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
        (ticker, view.get("date", ""), view.get("firm", ""), view.get("analyst"),
         view.get("rating"), view.get("prior_rating"),
         view.get("target_price"), view.get("prior_target"),
         view.get("currency", "USD"), view.get("key_thesis"),
         json.dumps(view.get("catalysts", []), ensure_ascii=False),
         json.dumps(view.get("risks", []), ensure_ascii=False),
         json.dumps(view.get("estimates", {}), ensure_ascii=False),
         source_file))
    conn.commit()
    conn.close()
    print(f"[OK] 观点已记录: {ticker} / {view.get('firm', 'N/A')}")


def load_views(ticker: str) -> list[dict]:
    conn = _get_views_db()
    rows = conn.execute(
        "SELECT * FROM views WHERE ticker = ? ORDER BY date DESC", (ticker,)
    ).fetchall()
    conn.close()
    cols = ["id", "ticker", "date", "firm", "analyst", "rating", "prior_rating",
            "target_price", "prior_target", "currency", "key_thesis", "catalysts",
            "risks", "estimates", "source_file", "created_at"]
    views = []
    for row in rows:
        v = dict(zip(cols, row))
        try:
            v["catalysts"] = json.loads(v["catalysts"]) if v["catalysts"] else []
            v["risks"] = json.loads(v["risks"]) if v["risks"] else []
            v["estimates"] = json.loads(v["estimates"]) if v["estimates"] else {}
        except (json.JSONDecodeError, TypeError):
            pass
        views.append(v)
    return views


def generate_view_summary(ticker: str) -> str:
    views = load_views(ticker)
    if not views:
        return ""
    lines = [f"# {ticker} 券商观点汇总", "", f"_最后更新: {datetime.now().strftime('%Y-%m-%d')}_", "",
             "## 当前共识", "", "| 券商 | 分析师 | 评级 | 目标价 | 最近更新 |",
             "|------|--------|------|--------|---------|"]
    for v in views:
        tp = v.get("target_price")
        curr = v.get("currency", "USD")
        tp_str = f"${tp}" if tp and curr == "USD" else f"{tp} {curr}" if tp else "N/A"
        lines.append(f"| {v.get('firm','N/A')} | {v.get('analyst','N/A')} | {v.get('rating','N/A')} | {tp_str} | {v.get('date','N/A')} |")

    ratings = [v.get("rating", "").lower() for v in views if v.get("rating")]
    buy = sum(1 for r in ratings if r in ("buy", "overweight"))
    hold = sum(1 for r in ratings if r in ("hold", "neutral", "equal-weight"))
    sell = sum(1 for r in ratings if r in ("sell", "underweight"))
    targets = [v["target_price"] for v in views if v.get("target_price")]
    avg = sum(targets) / len(targets) if targets else 0
    lines.extend(["", f"**共识分布:** {buy} Buy / {hold} Hold / {sell} Sell"])
    if avg:
        lines.append(f"**平均目标价:** ${avg:.0f}")
    return "\n".join(lines)


def display_history(ticker: str):
    views = load_views(ticker)
    if not views:
        print(f"No analyst views recorded for {ticker}")
        return
    print(f"\n{ticker} 券商观点变化")
    print("\u2501" * 60)
    print(f"{'日期':<12} {'券商':<12} {'评级变化':<20} {'目标价变化':<16}")
    print("\u2501" * 60)
    for v in views:
        date = v.get("date", "N/A")
        firm = str(v.get("firm", "N/A"))[:10]
        rating = v.get("rating", "N/A")
        prior_r = v.get("prior_rating")
        rating_str = f"{prior_r} -> {rating}" if prior_r else rating
        tp = v.get("target_price")
        prior_tp = v.get("prior_target")
        tp_str = f"${prior_tp} -> ${tp}" if prior_tp and tp else f"${tp}" if tp else "N/A"
        print(f"{date:<12} {firm:<12} {rating_str:<20} {tp_str:<16}")
    print("\u2501" * 60)
    ratings = [v.get("rating", "").lower() for v in views if v.get("rating")]
    buy = sum(1 for r in ratings if r in ("buy", "overweight"))
    hold = sum(1 for r in ratings if r in ("hold", "neutral", "equal-weight"))
    sell = sum(1 for r in ratings if r in ("sell", "underweight"))
    targets = [v["target_price"] for v in views if v.get("target_price")]
    avg = sum(targets) / len(targets) if targets else 0
    print(f"共识: {buy} Buy / {hold} Hold / {sell} Sell  平均目标价: ${avg:.0f}")


# ── QoQ Comparison ───────────────────────────────────────────────────────────

def generate_comparison(current: dict, previous: dict | None, fields: list[dict]) -> str:
    """Generate QoQ comparison markdown table."""
    if not previous:
        return "_首次提取，无历史数据可对比_"

    lines = [
        f"## {current['ticker']} {current['quarter']} vs {previous['quarter']}",
        "",
        "| 指标 | 上季 | 本季 | 变化 | 信号 |",
        "|------|------|------|------|------|",
    ]

    cur_m = current.get("metrics", {})
    prev_m = previous.get("metrics", {})

    for f in fields:
        key = f["key"]
        if f["type"] == "text":
            continue  # Skip text fields in comparison table

        cur_val = cur_m.get(key, {}).get("value") if key in cur_m else None
        prev_val = prev_m.get(key, {}).get("value") if key in prev_m else None

        if cur_val is None and prev_val is None:
            continue

        cur_str = str(cur_val) if cur_val is not None else "—"
        prev_str = str(prev_val) if prev_val is not None else "—"

        # Calculate delta for percent/number
        signal = ""
        delta_str = ""
        if cur_val is not None and prev_val is not None:
            try:
                if f["type"] == "percent":
                    c = float(str(cur_val).replace("%", "").strip())
                    p = float(str(prev_val).replace("%", "").strip())
                    d = c - p
                    delta_str = f"{d:+.1f}pp"
                    signal = "🟢 加速" if d > 2 else ("🔴 减速" if d < -2 else "—")
                elif f["type"] == "number":
                    c = float(str(cur_val).replace(",", ""))
                    p = float(str(prev_val).replace(",", ""))
                    if p != 0:
                        pct = (c - p) / abs(p) * 100
                        delta_str = f"{pct:+.1f}%"
                        signal = "🟢" if pct > 5 else ("🔴" if pct < -5 else "—")
            except (ValueError, TypeError):
                delta_str = "—"

        if cur_val is None:
            signal = "⚠️ 缺失"
        elif prev_val is None:
            signal = "🆕"

        lines.append(f"| {f['label']} | {prev_str} | {cur_str} | {delta_str} | {signal} |")

    return "\n".join(lines)


# ── Obsidian Output ──────────────────────────────────────────────────────────

def save_to_obsidian(current: dict, comparison: str, fields: list[dict], report_summary: str = ""):
    """Save extraction result, full summary, and comparison to Obsidian vault."""
    ticker = current["ticker"]
    quarter = current["quarter"]
    today = datetime.now().strftime("%Y-%m-%d")

    # Build note content — structure: summary first, then QoQ, then raw KPIs
    metrics = current.get("metrics", {})

    # Group metrics by segment for the KPI appendix
    segment_data = {}
    for f in fields:
        seg = f["segment"]
        if seg not in segment_data:
            segment_data[seg] = []
        m = metrics.get(f["key"], {})
        val = m.get("value") if m else None
        page = m.get("source_page") if m else None
        segment_data[seg].append({
            "label": f["label"],
            "value": val,
            "page": page,
        })

    body_parts = []

    # Part 1: Full report summary (the rich content)
    if report_summary:
        body_parts.append(report_summary.strip())
        body_parts.append("")

    # Part 2: QoQ comparison table
    body_parts.extend(["---", "", "## QoQ 对比", ""])
    body_parts.append(comparison)

    # Part 3: Raw KPI appendix (collapsible)
    body_parts.extend(["", "---", "", "<details>", "<summary>📊 KPI 提取明细（点击展开）</summary>", ""])
    for seg, items in segment_data.items():
        body_parts.append(f"#### {seg}")
        body_parts.append("| 指标 | 值 | 页码 |")
        body_parts.append("|------|-----|------|")
        for item in items:
            v = str(item["value"]) if item["value"] is not None else "—"
            p = str(item["page"]) if item["page"] is not None else "—"
            body_parts.append(f"| {item['label']} | {v} | p.{p} |")
        body_parts.append("")
    body_parts.append("</details>")

    body_parts.extend([
        "",
        "---",
        f"_来源: {current.get('source', 'N/A')} | 提取时间: {today}_",
    ])

    frontmatter = {
        "tags": ["卖方跟踪", ticker],
        "ticker": ticker,
        "quarter": quarter,
        "source": current.get("source", ""),
        "date": today,
    }

    fm_yaml = yaml.dump(frontmatter, allow_unicode=True, default_flow_style=False).strip()
    content = f"---\n{fm_yaml}\n---\n\n# {ticker} {quarter} 卖方跟踪\n\n" + "\n".join(body_parts)

    # Save (overwrite if exists)
    rel_path = f"研究/卖方跟踪/{ticker}/{today} {ticker} {quarter} 卖方跟踪.md"
    full_path = VAULT / rel_path
    full_path.parent.mkdir(parents=True, exist_ok=True)
    with open(full_path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"✅ Obsidian 笔记已保存: {rel_path}")

    try:
        update_dashboard()
        print("✅ Dashboard 已更新")
    except Exception as e:
        print(f"⚠️ Dashboard 更新失败: {e}")


# ── Main ─────────────────────────────────────────────────────────────────────

def detect_ticker_from_filename(filename: str) -> str | None:
    """Try to detect ticker from filename."""
    name = Path(filename).stem.lower()
    mapping = {
        "google": "GOOG", "alphabet": "GOOG",
        "meta": "META", "facebook": "META",
        "microsoft": "MSFT", "msft": "MSFT",
        "apple": "AAPL",
        "amazon": "AMZN", "amzn": "AMZN",
        "nvidia": "NVDA",
        "tesla": "TSLA",
    }
    for keyword, ticker in mapping.items():
        if keyword in name:
            return ticker
    return None


def detect_quarter_from_filename(filename: str) -> str | None:
    """Try to detect quarter from filename, e.g. '25Q4' -> 'Q4-2025'."""
    name = Path(filename).stem
    # Match patterns like 25Q4, Q4 2025, 2025Q4, etc.
    m = re.search(r"(\d{2,4})\s*Q(\d)", name, re.IGNORECASE)
    if m:
        year = m.group(1)
        q = m.group(2)
        if len(year) == 2:
            year = "20" + year
        return f"Q{q}-{year}"
    m = re.search(r"Q(\d)\s*(\d{2,4})", name, re.IGNORECASE)
    if m:
        q = m.group(1)
        year = m.group(2)
        if len(year) == 2:
            year = "20" + year
        return f"Q{q}-{year}"
    return None


def save_generic_to_obsidian(ticker: str, view: dict, report_summary: str, source_file: str = ""):
    """Save generic (no-template) extraction result to Obsidian vault."""
    today = datetime.now().strftime("%Y-%m-%d")
    firm = view.get("firm", "Unknown")
    report_date = view.get("date", today)

    catalysts = view.get("catalysts", [])
    risks = view.get("risks", [])
    estimates = view.get("estimates", {})

    cat_lines = "\n".join(f"- {c}" for c in catalysts) if catalysts else "_未提及_"
    risk_lines = "\n".join(f"- {r}" for r in risks) if risks else "_未提及_"

    est_lines = []
    if estimates:
        for k, v in estimates.items():
            if v:
                est_lines.append(f"- {k}: {v}")
    est_block = "\n".join(est_lines) if est_lines else "_未提及_"

    rating = view.get("rating", "N/A")
    prior_r = view.get("prior_rating")
    rating_display = f"{prior_r} -> {rating}" if prior_r else rating

    tp = view.get("target_price")
    prior_tp = view.get("prior_target")
    curr = view.get("currency", "USD")
    tp_display = f"{prior_tp} -> {tp} {curr}" if prior_tp and tp else f"{tp} {curr}" if tp else "N/A"

    body = f"""## 核心观点摘要

{report_summary.strip() if report_summary else "_AI摘要未生成_"}

---

## 评级 & 目标价

| 项目 | 值 |
|------|-----|
| 评级 | {rating_display} |
| 目标价 | {tp_display} |
| 分析师 | {view.get('analyst', 'N/A')} |
| 报告日期 | {report_date} |

## 核心论点

{view.get('key_thesis', '_未提取_')}

## 催化剂

{cat_lines}

## 风险

{risk_lines}

## 预估数据

{est_block}

---

_来源: {firm} | 提取时间: {today}_
"""

    frontmatter = {
        "tags": ["卖方跟踪", ticker],
        "ticker": ticker,
        "firm": firm,
        "rating": rating,
        "target_price": tp,
        "currency": curr,
        "report_date": report_date,
        "date": today,
    }
    fm_yaml = yaml.dump(frontmatter, allow_unicode=True, default_flow_style=False).strip()
    content = f"---\n{fm_yaml}\n---\n\n# {ticker} 卖方观点 — {firm} ({report_date})\n\n{body}"

    rel_path = f"研究/卖方跟踪/{ticker}/{today} {ticker} {firm} 卖方观点.md"
    full_path = VAULT / rel_path
    full_path.parent.mkdir(parents=True, exist_ok=True)
    with open(full_path, "w", encoding="utf-8") as fh:
        fh.write(content)
    print(f"[OK] Obsidian 笔记已保存: {rel_path}")

    # Update summary note
    summary_content = generate_view_summary(ticker)
    if summary_content:
        summary_path = VAULT / f"研究/卖方跟踪/{ticker}/{ticker} 券商观点汇总.md"
        with open(summary_path, "w", encoding="utf-8") as fh:
            fh.write(summary_content)
        print(f"[OK] 汇总笔记已更新: 研究/卖方跟踪/{ticker}/{ticker} 券商观点汇总.md")

    try:
        update_dashboard()
        print("[OK] Dashboard 已更新")
    except Exception as e:
        print(f"[WARN] Dashboard 更新失败: {e}")

    return full_path


def main():
    parser = argparse.ArgumentParser(description="Sellside Report Tracker")
    parser.add_argument("pdf", nargs="?", help="Path to sellside report (PDF or PPTX)")
    parser.add_argument("--ticker", help="Ticker symbol (auto-detected if omitted)")
    parser.add_argument("--quarter", help="Quarter like Q4-2025 (auto-detected if omitted)")
    parser.add_argument("--no-obsidian", action="store_true", help="Skip Obsidian save")
    parser.add_argument("--dry-run", action="store_true", help="Extract only, don't save")
    parser.add_argument("--history", metavar="TICKER", help="Display view history for a ticker")
    args = parser.parse_args()

    # ── History mode ──────────────────────────────────────────────────────────
    if args.history:
        display_history(args.history.upper())
        return

    # ── Require PDF arg when not in history mode ──────────────────────────────
    if not args.pdf:
        print("Error: PDF path required (or use --history TICKER)", file=sys.stderr)
        sys.exit(1)

    pdf_path = args.pdf
    if not os.path.exists(pdf_path):
        print(f"Error: File not found: {pdf_path}", file=sys.stderr)
        sys.exit(1)

    # Detect ticker & quarter
    ticker = args.ticker or detect_ticker_from_filename(pdf_path)
    quarter = args.quarter or detect_quarter_from_filename(pdf_path)

    if not ticker:
        print("Error: Cannot detect ticker. Use --ticker GOOG", file=sys.stderr)
        sys.exit(1)

    ticker = ticker.upper()
    print(f"[PDF] {Path(pdf_path).name}")
    print(f"[Ticker] {ticker}" + (f" | Quarter: {quarter}" if quarter else ""))

    # ── Try to load CSV template ───────────────────────────────────────────────
    fields = load_template(ticker)

    # ── Extract PDF text (common to both modes) ───────────────────────────────
    print("[...] 提取文本...")
    pdf_text = extract_file_text(pdf_path)
    page_count = pdf_text.count("[PAGE ")
    print(f"    {page_count} 页，{len(pdf_text)} 字符")

    # ══════════════════════════════════════════════════════════════════════════
    # GENERIC MODE — no CSV template found
    # ══════════════════════════════════════════════════════════════════════════
    if fields is None:
        if not quarter:
            quarter = "unknown"
        print(f"[INFO] 未找到 {ticker} 的 CSV 模板，启用通用提取模式")
        print("[...] 调用 Gemini: 通用结构化提取...")
        generic_prompt = build_generic_extraction_prompt(pdf_text, ticker)
        raw_json = call_gemini(generic_prompt)
        print(f"    [OK] 通用提取完成")

        # Parse JSON response — strip fences if present
        cleaned_json = raw_json.strip()
        if cleaned_json.startswith("```"):
            cleaned_json = re.sub(r"^```\w*\n?", "", cleaned_json)
            cleaned_json = re.sub(r"\n?```$", "", cleaned_json)

        try:
            view = json.loads(cleaned_json)
        except json.JSONDecodeError as e:
            print(f"[ERR] JSON 解析失败: {e}", file=sys.stderr)
            print("--- Raw response (first 2000 chars) ---")
            print(raw_json[:2000])
            sys.exit(1)

        print(f"    firm={view.get('firm','?')} | rating={view.get('rating','?')} | TP={view.get('target_price','?')}")

        if args.dry_run:
            print("\n[DRY RUN] 完成，未保存数据")
            print(json.dumps(view, ensure_ascii=False, indent=2))
            return

        # Save view to SQLite
        save_view(ticker, view, source_file=str(pdf_path))

        # Generate narrative summary
        print("[...] 调用 Gemini: 生成叙事摘要...")
        # Use a lightweight generic summary prompt
        summary_prompt = f"""你是一位资深买方分析师。请将以下卖方报告转化为简洁的投资笔记（中文，Markdown 格式）。

涵盖：核心观点、主要业务进展、估值逻辑、风险提示。
保留具体数字，标注来源页码如 (p.5)。

## 报告全文
{pdf_text}"""
        report_summary = call_gemini(summary_prompt)
        print(f"    [OK] 摘要完成 ({len(report_summary)} chars)")

        # Save to Obsidian + update summary note
        if not args.no_obsidian:
            note_path = save_generic_to_obsidian(ticker, view, report_summary, source_file=str(pdf_path))

        # Vector memory upsert
        try:
            sys.path.insert(0, str(Path.home() / ".claude" / "skills"))
            from shared.vector_memory import upsert_from_file
            if not args.no_obsidian:
                upsert_from_file(str(note_path))
                print("[OK] Vector memory 已更新")
        except Exception as e:
            print(f"[WARN] Vector memory 更新失败: {e}")

        print(f"\n[DONE] 通用提取完成 | {ticker} / {view.get('firm', 'N/A')}")
        return

    # ══════════════════════════════════════════════════════════════════════════
    # TEMPLATE MODE — CSV template found (existing flow, preserved unchanged)
    # ══════════════════════════════════════════════════════════════════════════
    if not quarter:
        print("Error: Cannot detect quarter. Use --quarter Q4-2025", file=sys.stderr)
        sys.exit(1)

    print(f"[Template] {len(fields)} 个指标 ({TEMPLATES_DIR / f'{ticker.lower()}.csv'})")

    # Build prompt & call AI — two calls: KPI extraction + full summary
    fields_text = template_to_prompt_fields(fields)
    prompt = build_extraction_prompt(pdf_text, fields_text, ticker, quarter)
    summary_prompt = build_summary_prompt(pdf_text, ticker, quarter)
    print("[...] 调用 Gemini: KPI 提取 + 全文摘要...")
    raw_response = call_gemini(prompt)
    print("    [OK] KPI 提取完成")
    print("    [...] 生成全文摘要中...")
    report_summary = call_gemini(summary_prompt)
    print(f"    [OK] 全文摘要完成 ({len(report_summary)} chars)")

    # Parse YAML response — strip markdown code fence if present
    cleaned = raw_response.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```\w*\n?", "", cleaned)
        cleaned = re.sub(r"\n?```$", "", cleaned)

    try:
        data = yaml.safe_load(cleaned)
    except yaml.YAMLError as e:
        print(f"[ERR] YAML 解析失败: {e}", file=sys.stderr)
        print("--- Raw response ---")
        print(raw_response[:2000])
        sys.exit(1)

    # Count extracted metrics
    metrics = data.get("metrics", {})
    filled = sum(1 for m in metrics.values() if m and m.get("value") is not None)
    total = len(fields)
    print(f"[OK] 提取完成: {filled}/{total} 指标已填充")

    # Load history & validate
    history = load_history(ticker)
    prev = get_previous_quarter(history, quarter)
    warnings = validate_extraction(data, fields, prev)

    if warnings:
        print(f"\n[WARN] {len(warnings)} 个警告:")
        for w in warnings:
            print(f"  {w}")

    # Print summary table
    print(f"\n{'─' * 60}")
    print(f"{'指标':<30} {'值':<20} {'页码':<5}")
    print(f"{'─' * 60}")
    for f in fields:
        key = f["key"]
        m = metrics.get(key, {})
        val = m.get("value", "—") if m else "—"
        page = m.get("source_page", "—") if m else "—"
        label = f["label"][:28]
        print(f"{label:<30} {str(val):<20} p.{page}")
    print(f"{'─' * 60}")

    if args.dry_run:
        print("\n[DRY RUN] 完成，未保存数据")
        return

    # Save to history
    history[quarter] = data
    save_history(ticker, history)

    # Generate comparison
    comparison = generate_comparison(data, prev, fields)
    print(f"\n{comparison}")

    # Save to Obsidian
    if not args.no_obsidian:
        save_to_obsidian(data, comparison, fields, report_summary)

    print(f"\n[DONE] 完成！{filled}/{total} 指标 | {len(warnings)} 警告")


if __name__ == "__main__":
    main()
